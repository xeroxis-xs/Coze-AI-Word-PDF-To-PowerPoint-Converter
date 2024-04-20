from runtime import Args
from typings.create_powerpoint_for_download.create_powerpoint_for_download import Input, Output

import uuid
import json
import requests
from PIL import Image
from io import BytesIO

from google.cloud import storage

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER_TYPE
from pptx.dml.color import RGBColor


"""
Each file needs to export a function named `handler`. This function is the entrance to the Tool.

Parameters:
args: parameters of the entry function.
args.input - input parameters, you can get test input value by args.input.xxx.
args.logger - logger instance used to print logs, injected by runtime.

Remember to fill in input/output in Metadata, it helps LLM to recognize and use tool.

Return:
The return data of the function, which should match the declared output parameters.
"""
class GCSClient:
    def __init__(self, key_dict, buket_name='coze-plugin-storage', source_file_name='powerpoint.pptx', destination_blob_name=str(uuid.uuid4()) + '.pptx'):
        self.storage_client = storage.Client.from_service_account_info(key_dict)
        self.bucket_name = buket_name
        self.source_file_name = source_file_name
        self.destination_blob_name = destination_blob_name

    def upload_file_to_gcs(self):
        """Upload a file to a Google Cloud Storage bucket."""
        bucket = self.storage_client.bucket(self.bucket_name)
        blob = bucket.blob(self.destination_blob_name)
        blob.upload_from_filename(self.source_file_name)

        print(f"File {self.source_file_name} uploaded to {self.destination_blob_name}.")

    def upload_byte_io_to_gcs(self, ppt_io):
        """Upload the BytesIO object to a Google Cloud Storage bucket."""
        bucket = self.storage_client.bucket(self.bucket_name)
        blob = bucket.blob(self.destination_blob_name)
        # Upload the BytesIO object to the blob
        blob.upload_from_file(ppt_io, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')


    def generate_download_link(self, expiration=3600):
        """Generate a signed URL for downloading a file from a Google Cloud Storage bucket."""
        bucket = self.storage_client.bucket(self.bucket_name)
        blob = bucket.blob(self.destination_blob_name)
        url = blob.generate_signed_url(
            version="v4",
            expiration=expiration,
            method="GET"
        )
        return url

class PresentationCreator:
    def __init__(self, data, args):
        self.args = args
        self.data = data
        self.presentation = Presentation()
        self.margin = Inches(1)


    def set_aspect_ratio(self, aspect_ratio):
        # Set aspect ratio
        self.presentation.slide_width = Inches(10 * aspect_ratio)
        self.presentation.slide_height = Inches(10)

    def add_data(self):
        i = 0
        for slide_data in self.data:
            if i == 0:
                # Add a title slide without content
                slide_layout = self.presentation.slide_layouts[0] # 0 is title slide, 1 is title and content
                slide = self.presentation.slides.add_slide(slide_layout)

                # self.args.logger.info(f"request from: {slide_data.background_img}")
                # Download the image from the URL
                response = requests.get(slide_data.background_img)
                image = Image.open(BytesIO(response.content))

                # Process the image
                image = self.resize_and_convert_image(image)

                # Save the image to a BytesIO object
                image_stream = BytesIO()
                image.save(image_stream, format='PNG')
                image_stream.seek(0)

                # Add a background image
                bg = slide.shapes.add_picture(image_stream, 0, 0, width=self.presentation.slide_width, height=self.presentation.slide_height)

                # Move it to the background
                slide.shapes._spTree.remove(bg._element)
                slide.shapes._spTree.insert(2, bg._element)

                # Add a title
                title = slide.shapes.title
                title.text = slide_data.title
            else:
                # Add a slide with a title and content layout
                slide_layout = self.presentation.slide_layouts[1]  # 0 is title slide, 1 is title and content
                slide = self.presentation.slides.add_slide(slide_layout)

                # Download the image from the URL
                response = requests.get(slide_data.background_img)
                image = Image.open(BytesIO(response.content))

                # Process the image
                image = self.resize_and_convert_image(image)

                # Save the image to a BytesIO object
                image_stream = BytesIO()
                image.save(image_stream, format='PNG')
                image_stream.seek(0)

                # Add a background image
                bg = slide.shapes.add_picture(image_stream, 0, 0, width=self.presentation.slide_width, height=self.presentation.slide_height)

                # Move it to the background
                slide.shapes._spTree.remove(bg._element)
                slide.shapes._spTree.insert(2, bg._element)

                # Add a title
                title = slide.shapes.title
                title.text = slide_data.title

                # Add multiple contents
                content_placeholder = slide.placeholders[1]

                for slide_data_content in slide_data.contents:
                    c = content_placeholder.text_frame.add_paragraph()
                    c.text = slide_data_content
            i += 1



    def set_size_and_position(self):
        i = 0
        for slide in self.presentation.slides:
            if i == 0:
                # title slide
                title = slide.shapes.title
                title.width = self.presentation.slide_width - 2 * self.margin
                title.height = Inches(10)
                title.left = self.margin
                title.top = 0

                # Delete subtitle placeholder
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.SUBTITLE:
                            sp = shape._element
                            sp.getparent().remove(sp)
            else:
                # title and content slide
                title = slide.shapes.title
                title.width = self.presentation.slide_width - 2 * self.margin
                title.height = Inches(1.5)
                title.left = self.margin
                title.top = self.margin

                # content placeholder
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.OBJECT:
                            shape.width = self.presentation.slide_width - 2 * self.margin
                            shape.left = self.margin
                            shape.top = title.top + title.height + Inches(0.2)  # add a 0.5 inch gap
                            # Calculate the maximum height of the content
                            max_content_height = self.presentation.slide_height - shape.top - self.margin
                            shape.height = max_content_height
            i += 1

    def set_background(self, background_color):
        for slide in self.presentation.slides:
            background = slide.background
            fill = background.fill
            fill.solid()
            fill.fore_color.rgb = self.hex_to_rgb(background_color)


    def set_title_slide_text_format(self, text_format):
        for slide in self.presentation.slides:
            title = slide.shapes.title
            for paragraph in title.text_frame.paragraphs:
                paragraph.font.size = Pt(text_format.title.font_size)
                paragraph.font.color.rgb = self.hex_to_rgb(text_format.title.font_color)
                paragraph.font.bold = text_format.title.bold
                paragraph.font.italic = text_format.title.italic
                paragraph.font.underline = text_format.title.underline
                paragraph.line_spacing = text_format.title.line_spacing
                paragraph.alignment = getattr(PP_ALIGN, text_format.title.alignment.upper())
            break;

    def set_content_slide_text_format(self, text_format):
        i = 0
        for slide in self.presentation.slides:
            if i != 0:
                for shape in slide.shapes:
                    if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
                        if shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.TITLE:
                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.size = Pt(text_format.title.font_size)
                                paragraph.font.color.rgb = self.hex_to_rgb(text_format.title.font_color)
                                paragraph.font.bold = text_format.title.bold
                                paragraph.font.italic = text_format.title.italic
                                paragraph.font.underline = text_format.title.underline
                                paragraph.line_spacing = text_format.title.line_spacing
                                paragraph.alignment = getattr(PP_ALIGN, text_format.title.alignment.upper())
                        elif shape.placeholder_format.type == PP_PLACEHOLDER_TYPE.OBJECT:
                            # Remove the first paragraph from the placeholder
                            if shape.text_frame.paragraphs:
                                shape.text_frame._element.remove(shape.text_frame.paragraphs[0]._element)

                            for paragraph in shape.text_frame.paragraphs:
                                paragraph.font.size = Pt(text_format.content.font_size)
                                paragraph.font.color.rgb = self.hex_to_rgb(text_format.content.font_color)
                                paragraph.font.bold = text_format.content.bold
                                paragraph.font.italic = text_format.content.italic
                                paragraph.font.underline = text_format.content.underline
                                paragraph.line_spacing = text_format.content.line_spacing
                                paragraph.alignment = getattr(PP_ALIGN, text_format.content.alignment.upper())
            i+=1


    def hex_to_rgb(self, hex_color):
        hex_color = hex_color.lstrip('#')
        int_color = int(hex_color, 16)
        red = (int_color >> 16) & 255
        green = (int_color >> 8) & 255
        blue = int_color & 255
        return RGBColor(red, green, blue)

    def resize_and_convert_image(self, image):
        # Compress the image
        image = image.resize((int(image.width / 2.5), int(image.height / 2.5)), Image.Resampling.LANCZOS)

        # Change the opacity of the image
        image = image.convert("RGBA")
        for x in range(image.width):
            for y in range(image.height):
                r, g, b, a = image.getpixel((x, y))
                image.putpixel((x, y), (r, g, b, int(a * 0.35)))  # Change 0.5 to the desired opacity

        return image

    def save_presentation(self):
        # Save the presentation
        self.presentation.save('powerpoint.pptx')

    def save_presenation_as_bytes_io(self):
        # Save the presentation as BytesIO object
        presentation_stream = BytesIO()
        self.presentation.save(presentation_stream)
        presentation_stream.seek(0)
        return presentation_stream


def clean_url(url):

    # Remove invalid \u0026 characters from the URL
    cleaned_url = url.replace("\\u0026", "&")
    
    return cleaned_url
    

def handler(args: Args[Input])->Output:
    try:
        data = args.input.data
        config = args.input.config
    except Exception as e:
        error_string = f"Error processing input arguments. Error: {e}"
        return {"download_link": error_string}

    try:
        for slide in data:
            # args.logger.info(f"before: {slide.background_img}")
            slide.background_img = clean_url(slide.background_img)
            # args.logger.info(f"after: {slide.background_img}")
    except Exception as e:
        error_string = f"Error cleaning url. Error: {e}"
        return {"download_link": error_string}

    try:
        presentation = PresentationCreator(data, args)
        presentation.set_aspect_ratio(config.aspect_ratio)
        presentation.add_data()
        presentation.set_background(config.background_color)
        presentation.set_title_slide_text_format(config.title_slide)
        presentation.set_content_slide_text_format(config.title_content_slide)
        presentation.set_size_and_position()
        ppt_io = presentation.save_presenation_as_bytes_io()
    except Exception as e:
        error_string = f"Error creating presentation. Error: {e}"
        return {"download_link": error_string}


    key_dict = { "your key dict here" }


    try: 
        # Authenticate with the key file
        gcs_client = GCSClient(key_dict)

        gcs_client.upload_byte_io_to_gcs(ppt_io)

        download_link = gcs_client.generate_download_link()
    except Exception as e:
        error_string = f"Error Uploading created presentation deck on GCS. Error: {e}"
        return {"download_link": error_string}


    return {"download_link": download_link}
